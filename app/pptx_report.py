"""Generate a per-campaign PowerPoint report (the "📥 PowerPoint" button).

Built ON TOP of the agency's own deck design: app/assets/report_template.pptx
is the real P2026-027 PAO file with its slides stripped out but every layout,
master, background art and embedded font kept. Slides are recreated at the
exact coordinates measured from the original file, so the output matches the
team's format:

  1. cover  — "Influencer Report / Campaign : <name>"   (layout 1 artwork)
  2. group divider slides when the roster has >1 big group (layout 1)
  3. one slide per post (layout 2): platform logo + KOL name inside the header
     band, a composed post screenshot (header + caption + image) in the
     original picture box, link strip at the bottom, and the KPI/stat table
     with the original teal header + gray grid. Blank where data can't be
     fetched; KPI always blank. A KOL's platforms are consecutive slides.
  4. closing "Thank You" slide (layout 1)
"""
from __future__ import annotations

import hashlib
import io
import logging
import pathlib
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Emu, Pt
from sqlalchemy import select

from app.db import session_scope
from app.models import Campaign, ImageCache, ReportKol, ReportPost
from app.report_refresh import kol_links

log = logging.getLogger("pptx_report")

ASSETS = pathlib.Path(__file__).resolve().parent / "assets"
TEMPLATE = ASSETS / "report_template.pptx"
LOGO_DIR = ASSETS / "logos"
FONT = "Leelawadee"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
TEAL = RGBColor(0x0D, 0x9B, 0xB3)     # KPI header cell (from the original deck)
CELL_GRAY = RGBColor(0xE0, 0xE6, 0xEC)

PLAT_LABEL = {"tiktok": "TikTok", "facebook": "Facebook", "instagram": "Instagram",
              "youtube": "YouTube", "x": "X (Twitter)", "line": "LINE",
              "website": "Website", "other": "Link"}
PLAT_ORDER = {"tiktok": 0, "facebook": 1, "instagram": 2, "youtube": 3,
              "x": 4, "line": 5, "website": 6, "other": 7}
# fallback letter-circle colors for platforms with no logo asset
PLAT_FALLBACK = {"x": ("X", (17, 17, 17)), "line": ("L", (6, 199, 85)),
                 "website": ("W", (14, 165, 233)), "other": ("•", (100, 116, 139))}


def _fmt(n) -> str:
    return f"{int(n):,}" if n else ""


def _stat_rows(platform: str, p: Optional[ReportPost]) -> list:
    """(label, value) rows exactly as the team lays them out per platform —
    blank where the scraper has no data."""
    if platform == "facebook":
        eng = (p.likes + p.comments + p.shares) if p else 0
        return [("Reach", ""), ("View", _fmt(p.views) if p else ""),
                ("Engagement", _fmt(eng) if p else ""),
                ("Reactions", _fmt(p.likes) if p else ""),
                ("Comments", _fmt(p.comments) if p else ""),
                ("Shares", _fmt(p.shares) if p else "")]
    if platform == "instagram":
        return [("View", _fmt(p.views) if p else ""), ("Like", _fmt(p.likes) if p else ""),
                ("Comments", _fmt(p.comments) if p else ""), ("Shares", "")]
    if platform == "youtube":
        return [("View", _fmt(p.views) if p else ""), ("Like", _fmt(p.likes) if p else ""),
                ("Comments", _fmt(p.comments) if p else "")]
    if platform == "x":
        return [("View", _fmt(p.views) if p else ""), ("Like", _fmt(p.likes) if p else ""),
                ("Comments", _fmt(p.comments) if p else ""),
                ("Reposts", _fmt(p.shares) if p else "")]
    return [("Impression", ""), ("View", _fmt(p.views) if p else ""), ("Reach", ""),
            ("Like", _fmt(p.likes) if p else ""), ("Comment", _fmt(p.comments) if p else ""),
            ("Share", _fmt(p.shares) if p else ""), ("Save", _fmt(p.saves) if p else "")]


# ---------------------------------------------------------------------------
# fonts + images
# ---------------------------------------------------------------------------

_THAI_TTF = pathlib.Path("/tmp/NotoSansThai.ttf")
_THAI_URL = ("https://github.com/google/fonts/raw/main/ofl/notosansthai/"
             "NotoSansThai%5Bwdth%2Cwght%5D.ttf")


def _thai_font(size: int):
    """A Thai-capable font for the composed post card. Downloaded once (the
    server has no Thai system fonts); falls back to PIL's built-in font."""
    from PIL import ImageFont
    try:
        if not _THAI_TTF.exists():
            import httpx
            r = httpx.get(_THAI_URL, timeout=30, follow_redirects=True)
            if r.status_code == 200 and len(r.content) > 10000:
                _THAI_TTF.write_bytes(r.content)
        if _THAI_TTF.exists():
            return ImageFont.truetype(str(_THAI_TTF), size)
    except Exception:  # noqa: BLE001
        pass
    try:
        return ImageFont.load_default(size=size)
    except TypeError:  # Pillow < 10.1
        return ImageFont.load_default()


_UA = ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
       "(KHTML, like Gecko) Chrome/125.0.0.0 Safari/537.36")


def _referer_for(url: str) -> Optional[str]:
    """CDNs hotlink-check against their own site — a foreign Referer can 403."""
    u = url.lower()
    if "tiktok" in u:
        return "https://www.tiktok.com/"
    if "fbcdn" in u or "facebook" in u:
        return "https://www.facebook.com/"
    if "instagram" in u or "cdninstagram" in u:
        return "https://www.instagram.com/"
    return None


def _image_bytes(session, url: Optional[str]) -> Optional[bytes]:
    if not url:
        return None
    h = hashlib.sha256(url.encode("utf-8")).hexdigest()[:40]
    row = session.get(ImageCache, h)
    if row and row.data:
        return row.data
    import httpx
    for ref in (_referer_for(url), None):  # matching referer, then bare retry
        try:
            headers = {"User-Agent": _UA,
                       "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8"}
            if ref:
                headers["Referer"] = ref
            r = httpx.get(url, timeout=8, follow_redirects=True, headers=headers)
            ct = (r.headers.get("content-type") or "").lower()
            if r.status_code == 200 and r.content and ct.startswith("image/"):
                return r.content
        except Exception:  # noqa: BLE001
            continue
    return None


def _platform_logo(plat: str) -> Optional[bytes]:
    """Logo PNG for the slide's top-left corner. TikTok/FB/IG/YT logos are the
    ones extracted from the team's original deck; others get a letter circle."""
    f = LOGO_DIR / f"{plat}.png"
    if f.exists():
        try:
            return f.read_bytes()
        except Exception:  # noqa: BLE001
            pass
    letter, rgb = PLAT_FALLBACK.get(plat, PLAT_FALLBACK["other"])
    try:
        from PIL import Image, ImageDraw
        side = 256
        im = Image.new("RGBA", (side, side), (0, 0, 0, 0))
        d = ImageDraw.Draw(im)
        d.ellipse((0, 0, side, side), fill=rgb + (255,))
        fnt = _thai_font(140)
        bbox = d.textbbox((0, 0), letter, font=fnt)
        d.text(((side - bbox[2] + bbox[0]) / 2 - bbox[0],
                (side - bbox[3] + bbox[1]) / 2 - bbox[1]),
               letter, font=fnt, fill=(255, 255, 255, 255))
        out = io.BytesIO()
        im.save(out, "PNG")
        return out.getvalue()
    except Exception:  # noqa: BLE001
        return None


def _fmt_short(n) -> str:
    n = int(n or 0)
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M".replace(".0M", "M")
    if n >= 1_000:
        return f"{n/1_000:.1f}K".replace(".0K", "K")
    return str(n)


def _wrap(draw, text, font, max_w, max_lines):
    """Greedy character wrap (Thai has no spaces); ellipsis past max_lines."""
    lines, cur = [], ""
    for ch in (text or "").replace("\r", ""):
        if ch == "\n":
            lines.append(cur)
            cur = ""
        elif draw.textlength(cur + ch, font=font) > max_w:
            lines.append(cur)
            cur = ch
        else:
            cur += ch
        if len(lines) >= max_lines:
            break
    if len(lines) < max_lines and cur:
        lines.append(cur)
    if len(lines) >= max_lines:
        lines = lines[:max_lines]
        lines[-1] = (lines[-1][:-2] + "…") if len(lines[-1]) > 2 else lines[-1] + "…"
    return lines


def _circle_img(img_bytes: bytes, side: int):
    """Center-crop an image to a circle (profile pictures)."""
    from PIL import Image, ImageDraw
    im = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    m = min(im.size)
    im = im.crop(((im.width - m) // 2, (im.height - m) // 2,
                  (im.width + m) // 2, (im.height + m) // 2)).resize((side, side))
    mask = Image.new("L", (side, side), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, side, side), fill=255)
    im.putalpha(mask)
    return im


def _icon_heart(d, cx, cy, s, fill):
    r = s * 0.30
    d.ellipse((cx - s * 0.5, cy - s * 0.42, cx - s * 0.5 + 2 * r, cy - s * 0.42 + 2 * r), fill=fill)
    d.ellipse((cx + s * 0.5 - 2 * r, cy - s * 0.42, cx + s * 0.5, cy - s * 0.42 + 2 * r), fill=fill)
    d.polygon([(cx - s * 0.5, cy - s * 0.05), (cx + s * 0.5, cy - s * 0.05),
               (cx, cy + s * 0.52)], fill=fill)
    d.rectangle((cx - s * 0.5, cy - s * 0.18, cx + s * 0.5, cy - s * 0.02), fill=fill)


def _icon_comment(d, cx, cy, s, fill):
    d.ellipse((cx - s * 0.5, cy - s * 0.42, cx + s * 0.5, cy + s * 0.3), fill=fill)
    d.polygon([(cx - s * 0.25, cy + s * 0.2), (cx - s * 0.02, cy + s * 0.52),
               (cx + 0.12 * s, cy + s * 0.2)], fill=fill)


def _icon_bookmark(d, cx, cy, s, fill):
    d.polygon([(cx - s * 0.32, cy - s * 0.5), (cx + s * 0.32, cy - s * 0.5),
               (cx + s * 0.32, cy + s * 0.5), (cx, cy + s * 0.22),
               (cx - s * 0.32, cy + s * 0.5)], fill=fill)


def _icon_share(d, cx, cy, s, fill):
    d.polygon([(cx - s * 0.5, cy + s * 0.45), (cx + s * 0.05, cy - s * 0.05),
               (cx + s * 0.05, cy + s * 0.18)], fill=fill)
    d.polygon([(cx + s * 0.05, cy - s * 0.35), (cx + s * 0.5, cy),
               (cx + s * 0.05, cy + s * 0.35)], fill=fill)
    d.rectangle((cx - s * 0.12, cy - s * 0.12, cx + s * 0.12, cy + s * 0.12), fill=fill)


def _icon_thumb(d, cx, cy, s, circle, fill):
    d.ellipse((cx - s * 0.5, cy - s * 0.5, cx + s * 0.5, cy + s * 0.5), fill=circle)
    d.rectangle((cx - s * 0.28, cy - s * 0.02, cx - s * 0.12, cy + s * 0.26), fill=fill)
    d.rectangle((cx - s * 0.10, cy - s * 0.05, cx + s * 0.26, cy + s * 0.26), fill=fill)
    d.rectangle((cx - s * 0.04, cy - s * 0.28, cx + s * 0.10, cy + s * 0.02), fill=fill)


def _tiktok_card(cover: Optional[bytes], avatar: Optional[bytes], username: str,
                 date_s: str, caption: Optional[str], likes, comments, saves,
                 shares) -> Optional[io.BytesIO]:
    """Mock TikTok screenshot: full-bleed video frame, right action rail with
    like/comment/save/share counts, username + caption bottom-left."""
    try:
        from PIL import Image, ImageDraw
        W, H = 720, 1250
        card = Image.new("RGB", (W, H), (18, 18, 18))
        if cover:
            try:
                im = Image.open(io.BytesIO(cover)).convert("RGB")
                scale = max(W / im.width, H / im.height)
                im = im.resize((int(im.width * scale) or W, int(im.height * scale) or H))
                card.paste(im, ((W - im.width) // 2, (H - im.height) // 2))
            except Exception:  # noqa: BLE001
                pass
        # bottom gradient for text legibility
        grad = Image.new("L", (1, 320), 0)
        for i in range(320):
            grad.putpixel((0, i), int(190 * (i / 320)))
        shade = Image.new("RGB", (W, 320), (0, 0, 0))
        card.paste(shade, (0, H - 320), grad.resize((W, 320)))
        d = ImageDraw.Draw(card)

        # right action rail
        rail_x = W - 76
        y = int(H * 0.40)
        wht = (255, 255, 255)
        if avatar:
            try:
                av = _circle_img(avatar, 92)
                card.paste(av, (rail_x - 46, y - 46), av)
                d.ellipse((rail_x - 48, y - 48, rail_x + 48, y + 48),
                          outline=wht, width=3)
            except Exception:  # noqa: BLE001
                pass
        y += 128
        f_cnt = _thai_font(26)
        for icon, val in ((_icon_heart, likes), (_icon_comment, comments),
                          (_icon_bookmark, saves), (_icon_share, shares)):
            icon(d, rail_x, y, 56, wht)
            t = _fmt_short(val)
            d.text((rail_x - d.textlength(t, font=f_cnt) / 2, y + 40), t,
                   font=f_cnt, fill=wht)
            y += 128

        # bottom-left: username · date + caption
        f_user = _thai_font(30)
        f_cap = _thai_font(25)
        tx, ty = 28, H - 218
        d.text((tx, ty), f"{username} · {date_s}", font=f_user, fill=wht,
               stroke_width=1, stroke_fill=wht)
        ty += 46
        for lnn in _wrap(d, caption or "", f_cap, W - 150, 3):
            d.text((tx, ty), lnn, font=f_cap, fill=wht)
            ty += 34

        out = io.BytesIO()
        card.save(out, "PNG")
        out.seek(0)
        return out
    except Exception:  # noqa: BLE001
        log.exception("tiktok card failed")
        return io.BytesIO(cover) if cover else None


def _fb_card(avatar: Optional[bytes], page: str, date_s: str,
             caption: Optional[str], media: Optional[bytes], likes, comments,
             shares) -> Optional[io.BytesIO]:
    """Mock Facebook post: page header, full caption, media, engagement bar."""
    if not media and not caption:
        return None
    try:
        from PIL import Image, ImageDraw
        W, PAD = 760, 24
        f_name = _thai_font(28)
        f_date = _thai_font(21)
        f_cap = _thai_font(25)
        f_cnt = _thai_font(23)
        probe = ImageDraw.Draw(Image.new("RGB", (8, 8)))
        cap_lines = _wrap(probe, caption or "", f_cap, W - 2 * PAD, 8)
        head_h = 104
        cap_h = (len(cap_lines) * 34 + 10) if cap_lines else 0
        m_img = None
        if media:
            try:
                m_img = Image.open(io.BytesIO(media)).convert("RGB")
                mh = min(int(m_img.height * W / m_img.width), 860)
                m_img = m_img.resize((W, mh))
            except Exception:  # noqa: BLE001
                m_img = None
        media_h = m_img.height if m_img else 0
        foot_h = 64
        H = head_h + cap_h + media_h + foot_h

        card = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(card)
        d.rectangle((0, 0, W - 1, H - 1), outline=(224, 228, 234), width=2)

        # header: avatar circle + page name + date
        x = PAD
        if avatar:
            try:
                av = _circle_img(avatar, 60)
                card.paste(av, (PAD, 22), av)
                x = PAD + 60 + 14
            except Exception:  # noqa: BLE001
                pass
        d.text((x, 22), page, font=f_name, fill=(5, 5, 5))
        d.text((x, 60), date_s, font=f_date, fill=(101, 103, 107))

        y = head_h
        for lnn in cap_lines:
            d.text((PAD, y), lnn, font=f_cap, fill=(5, 5, 5))
            y += 34
        if cap_lines:
            y += 10
        if m_img:
            card.paste(m_img, (0, y))
            y += m_img.height

        # engagement bar
        d.line((PAD, y + 6, W - PAD, y + 6), fill=(224, 228, 234), width=2)
        cy = y + foot_h // 2 + 4
        _icon_thumb(d, PAD + 16, cy, 34, (24, 119, 242), (255, 255, 255))
        d.text((PAD + 42, cy - 14), _fmt_short(likes), font=f_cnt, fill=(101, 103, 107))
        right = f"{_fmt_short(comments)} comments · {_fmt_short(shares)} shares"
        d.text((W - PAD - d.textlength(right, font=f_cnt), cy - 14), right,
               font=f_cnt, fill=(101, 103, 107))

        out = io.BytesIO()
        card.save(out, "PNG")
        out.seek(0)
        return out
    except Exception:  # noqa: BLE001
        log.exception("fb card failed")
        return io.BytesIO(media) if media else None


def _compose_post_card(logo: Optional[bytes], name: str, caption: Optional[str],
                       img: Optional[bytes]) -> Optional[io.BytesIO]:
    """Generic post card (header + caption + media) for platforms without a
    dedicated mock (IG/YT/X/website)."""
    if not img and not caption:
        return None
    try:
        from PIL import Image, ImageDraw
        W, PAD = 760, 26
        f_name = _thai_font(30)
        f_cap = _thai_font(24)
        probe = ImageDraw.Draw(Image.new("RGB", (8, 8)))
        cap_lines = _wrap(probe, caption or "", f_cap, W - 2 * PAD, 6)
        head_h = 92
        cap_h = (len(cap_lines) * 32 + 14) if cap_lines else 0
        media = None
        if img:
            media = Image.open(io.BytesIO(img)).convert("RGB")
            mw = W - 2 * PAD
            mh = min(int(media.height * mw / media.width), 900)
            media = media.resize((mw, mh))
        media_h = (media.height + PAD) if media else 0
        H = head_h + cap_h + media_h + PAD

        card = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(card)
        d.rectangle((0, 0, W - 1, H - 1), outline=(229, 231, 235), width=2)
        x = PAD
        if logo:
            try:
                lg = Image.open(io.BytesIO(logo)).convert("RGBA").resize((52, 52))
                card.paste(lg, (PAD, 22), lg)
                x = PAD + 52 + 14
            except Exception:  # noqa: BLE001
                pass
        d.text((x, 30), name, font=f_name, fill=(15, 23, 42))
        y = head_h
        for lnn in cap_lines:
            d.text((PAD, y), lnn, font=f_cap, fill=(51, 65, 85))
            y += 32
        if cap_lines:
            y += 14
        if media:
            card.paste(media, (PAD, y))
        out = io.BytesIO()
        card.save(out, "PNG")
        out.seek(0)
        return out
    except Exception:  # noqa: BLE001
        log.exception("compose post card failed")
        if img:
            return io.BytesIO(img)
        return None


# ---------------------------------------------------------------------------
# slide helpers (coordinates measured from the original deck)
# ---------------------------------------------------------------------------

def _layout(prs, n: int):
    for lay in prs.slide_layouts:
        if str(lay.part.partname).endswith(f"slideLayout{n}.xml"):
            return lay
    return prs.slide_layouts[0]


def _add(prs, layout):
    s = prs.slides.add_slide(layout)
    for shp in list(s.shapes):  # drop cloned layout placeholders (empty boxes)
        if shp.is_placeholder:
            shp._element.getparent().remove(shp._element)
    return s


def _txt(slide, x, y, w, h, text, size, *, bold=False, color=BLACK, link=None,
         middle=False):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    if middle:  # original boxes use anchor="ctr" so text sits in the band
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT
    if link:
        try:
            run.hyperlink.address = link
        except Exception:  # noqa: BLE001
            pass
    return box


def _fit_picture(slide, fp, bx, by, bw, bh):
    """Place an image inside the original deck's picture box (fit + center)."""
    try:
        from PIL import Image
        data = fp.getvalue() if isinstance(fp, io.BytesIO) else fp
        iw, ih = Image.open(io.BytesIO(data)).size
        scale = min(bw / iw, bh / ih)
        w, h = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(io.BytesIO(data), Emu(bx + (bw - w) // 2),
                                 Emu(by + (bh - h) // 2), Emu(w), Emu(h))
    except Exception:  # noqa: BLE001
        pass


_LN_XML = ('<a:ln{side} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
           'w="12700" cap="flat" cmpd="sng" algn="ctr">'
           '<a:solidFill><a:srgbClr val="E0E6EC"/></a:solidFill>'
           '<a:prstDash val="solid"/></a:ln{side}>')


def _cell_borders(cell):
    """1pt E0E6EC grid lines on all sides — same as the original deck's table."""
    tcPr = cell._tc.get_or_add_tcPr()
    for i, side in enumerate(("L", "R", "T", "B")):
        tcPr.insert(i, parse_xml(_LN_XML.format(side=side)))


def _stat_table(slide, rows: list):
    """KPI + stat table at the original position/format: teal KPI header cell,
    gray blank KPI value, bold labels, gray 1pt grid, 0.5-inch centered rows."""
    n = len(rows) + 1
    x, y, w = 7360297, 1686539, 3883100
    row_h = 457200
    gf = slide.shapes.add_table(n, 2, Emu(x), Emu(y), Emu(w), Emu(row_h * n))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    tbl.columns[0].width = Emu(w // 2)
    tbl.columns[1].width = Emu(w - w // 2)
    for i in range(n):
        tbl.rows[i].height = Emu(row_h)

    def set_cell(cell, text, *, bold, fill, color=BLACK, size=14):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Emu(68700)
        cell.margin_right = Emu(68700)
        cell.margin_top = Emu(37775)
        cell.margin_bottom = Emu(37775)
        _cell_borders(cell)
        tf = cell.text_frame
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = text
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = FONT

    set_cell(tbl.cell(0, 0), "KPI", bold=True, fill=TEAL, color=WHITE, size=16)
    set_cell(tbl.cell(0, 1), "", bold=False, fill=CELL_GRAY)
    for i, (label, val) in enumerate(rows, start=1):
        set_cell(tbl.cell(i, 0), label, bold=True, fill=WHITE)
        set_cell(tbl.cell(i, 1), val, bold=False, fill=WHITE)


# ---------------------------------------------------------------------------
# deck assembly
# ---------------------------------------------------------------------------

def build_pptx(campaign_key: str) -> tuple[io.BytesIO, str]:
    """Build the deck from the agency template; returns (bytes, filename)."""
    with session_scope() as session:
        camp = session.get(Campaign, campaign_key)
        camp_name = camp.name if camp else campaign_key
        kols = session.scalars(
            select(ReportKol)
            .where(ReportKol.active.is_(True), ReportKol.campaign == campaign_key)
            .order_by(ReportKol.sort_order, ReportKol.id)  # source-file order
        ).all()
        posts_by = {}
        for p in session.scalars(select(ReportPost).where(
                ReportPost.campaign == campaign_key)).all():
            key = (p.username.lower(), p.platform or "tiktok")
            if key not in posts_by or p.views > posts_by[key].views:
                posts_by[key] = p

        prs = Presentation(str(TEMPLATE))
        lay_dark = _layout(prs, 1)   # cover / divider / thank-you artwork
        lay_stat = _layout(prs, 2)   # per-post artwork

        # ---- 1. cover (original slide 1 geometry) ----------------------------
        s = _add(prs, lay_dark)
        _txt(s, 4410075, 2383644, 7124700, 1397781,
             "Influencer Report", 44, color=WHITE, middle=True)
        _txt(s, 4410075, 3781425, 7124700, 640812,
             f"Campaign : {camp_name}", 18, color=WHITE, middle=True)

        # ---- 2-3. dividers + per-post slides ---------------------------------
        groups: dict[str, list[ReportKol]] = {}
        for k in kols:
            groups.setdefault(k.content_group or "KOL", []).append(k)
        show_dividers = len(groups) > 1

        for gname, gkols in groups.items():
            if show_dividers:  # original slide 9 geometry
                s = _add(prs, lay_dark)
                _txt(s, 4914900, 2469369, 5981700, 1397781, gname, 44,
                     color=WHITE, middle=True)
                _txt(s, 4914900, 3867150, 5981700, 640812,
                     f"{len(gkols)} KOLs", 18, color=WHITE, middle=True)

            for k in gkols:
                links = kol_links(k) or [{"platform": "tiktok", "url": k.url or "",
                                          "handle": k.username.lower()}]
                links = sorted(links, key=lambda l: PLAT_ORDER.get(l["platform"], 9))
                avatar = _image_bytes(session, k.avatar_url)
                for ln in links:
                    plat = ln["platform"]
                    p = posts_by.get((k.username.lower(), plat))
                    s = _add(prs, lay_stat)

                    # platform logo top-left (as in the original deck)
                    logo = _platform_logo(plat)
                    if logo:
                        s.shapes.add_picture(io.BytesIO(logo), Emu(187587),
                                             Emu(160544), Emu(879810), Emu(879810))

                    # KOL name + date inside the header band (anchor middle)
                    name = k.display or k.username
                    _txt(s, 1067397, 161170, 5789644, 662150, name, 20, middle=True)
                    posted = (p.posted_at.strftime("%d %b %Y")
                              if p and p.posted_at else "")
                    _txt(s, 1067397, 519577, 5681886, 536575,
                         f"Post Date: {posted}", 11, middle=True)

                    # composed post screenshot, styled per platform
                    shot = _image_bytes(session, (p.cover_url if p else None))
                    cap = p.caption if p else None
                    if plat == "tiktok":
                        dt_s = (f"{p.posted_at.month}-{p.posted_at.day}"
                                if p and p.posted_at else "")
                        card = _tiktok_card(shot, avatar, k.username, dt_s, cap,
                                            p.likes if p else 0, p.comments if p else 0,
                                            p.saves if p else 0, p.shares if p else 0)
                    elif plat == "facebook":
                        card = _fb_card(avatar, name, posted, cap, shot,
                                        p.likes if p else 0, p.comments if p else 0,
                                        p.shares if p else 0)
                    else:
                        card = _compose_post_card(logo, name, cap, shot)
                    if card:
                        _fit_picture(s, card, 908482, 1349012, 2628527, 4569755)

                    _stat_table(s, _stat_rows(plat, p))

                    if ln.get("url"):
                        _txt(s, 187587, 6277800, 10372464, 580200,
                             f"Link : {ln['url']}", 12, link=ln["url"])

        # ---- 4. thank you (original slide 16 geometry) ------------------------
        s = _add(prs, lay_dark)
        _txt(s, 5972175, 2621769, 4219575, 1397781, "Thank You", 44,
             color=WHITE, middle=True)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf, f"{camp_name} - Influencer Report.pptx"
